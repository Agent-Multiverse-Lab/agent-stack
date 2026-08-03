from __future__ import annotations

import json
import tempfile
import unittest
from pathlib import Path
from unittest import mock

from src.knowledge.extractor import ExtractorResult
from src.knowledge.parser import (
    DocParser,
    DocxParser,
    HtmlParser,
    ImageParser,
    MarkdownParser,
    PdfParser,
    PptxParser,
    TableParser,
    TextParser,
    get_parser,
    resolve_parser,
)
from src.knowledge.parser.pdf_parser import PdfParseError, PlainPdfParser


def _small_text_pdf() -> bytes:
    stream = b"BT /F1 12 Tf 72 720 Td (Hello PDF) Tj ET"
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>"
        ),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length %d >>\nstream\n" % len(stream) + stream + b"\nendstream",
    ]
    return _build_pdf(objects)


def _build_pdf(objects: list[bytes]) -> bytes:
    data = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for number, obj in enumerate(objects, start=1):
        offsets.append(len(data))
        data.extend(f"{number} 0 obj\n".encode() + obj + b"\nendobj\n")
    xref = len(data)
    data.extend(
        f"xref\n0 {len(objects) + 1}\n0000000000 65535 f \n".encode()
    )
    for offset in offsets[1:]:
        data.extend(f"{offset:010d} 00000 n \n".encode())
    data.extend(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref}\n%%EOF\n"
        ).encode()
    )
    return bytes(data)


class ParserTest(unittest.IsolatedAsyncioTestCase):
    def test_registry_routes_without_guessing(self) -> None:
        self.assertIs(resolve_parser(filename="sample.pdf"), PdfParser)
        self.assertIsInstance(get_parser(filename="sample.xlsx"), TableParser)
        self.assertIs(resolve_parser(filename="sample.csv"), TableParser)
        self.assertIs(resolve_parser(content_type="image/png"), ImageParser)
        self.assertIs(resolve_parser(parser_type="markdown"), MarkdownParser)
        with self.assertRaises(ValueError):
            resolve_parser(filename="sample.pdf", parser_type="csv")
        with self.assertRaises(ValueError):
            resolve_parser()

    async def test_text_csv_markdown_and_html(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            text_path = root / "sample.txt"
            text_path.write_text("第一段\n\n第二段", encoding="utf-8")
            csv_path = root / "sample.csv"
            csv_path.write_text("name,value\nA,1\n", encoding="utf-8")
            md_path = root / "sample.md"
            md_path.write_text("# 标题\n\n正文", encoding="utf-8")
            html_path = root / "sample.html"
            html_path.write_text(
                "<html><body><h1>标题</h1><p>正文</p><script>x</script></body></html>",
                encoding="utf-8",
            )

            self.assertIn("第一段", await TextParser().parse(text_path))
            self.assertIn("| name | value |", await TableParser().parse(csv_path))
            self.assertEqual("# 标题\n\n正文", await MarkdownParser().parse(md_path))
            self.assertNotIn("script", await HtmlParser().parse(html_path))

            for parser, path in (
                (TextParser(), text_path),
                (TableParser(), csv_path),
                (MarkdownParser(), md_path),
                (HtmlParser(), html_path),
            ):
                json.dumps(await parser.parse(path, as_json=True), ensure_ascii=False)

    async def test_office_parsers(self) -> None:
        from docx import Document
        from openpyxl import Workbook
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)

            docx_path = root / "sample.docx"
            document = Document()
            document.add_heading("标题", level=1)
            document.add_paragraph("正文")
            document.save(docx_path)

            xlsx_path = root / "sample.xlsx"
            workbook = Workbook()
            worksheet = workbook.active
            worksheet.title = "数据"
            worksheet.append(["name", "value"])
            worksheet.append(["A", 1])
            workbook.save(xlsx_path)
            workbook.close()

            pptx_path = root / "sample.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "标题"
            slide.placeholders[1].text = "正文"
            presentation.save(pptx_path)

            self.assertIn("# 标题", await DocxParser().parse(docx_path))
            self.assertIn("## 数据", await TableParser().parse(xlsx_path))
            self.assertIn("幻灯片 1", await PptxParser().parse(pptx_path))

            json.dumps(
                await DocxParser().parse(docx_path, as_json=True),
                ensure_ascii=False,
            )
            json.dumps(
                await TableParser().parse(xlsx_path, as_json=True),
                ensure_ascii=False,
            )
            json.dumps(
                await PptxParser().parse(pptx_path, as_json=True),
                ensure_ascii=False,
            )

    async def test_pdf_ocr_falls_back_to_plain(self) -> None:
        calls: list[str] = []

        class StubPdf:
            def __init__(
                self,
                name: str,
                result: str | dict[str, object] | None = None,
                error: Exception | None = None,
            ) -> None:
                self.name = name
                self.result = result
                self.error = error

            async def parse(self, *_: object, **__: object) -> object:
                calls.append(self.name)
                if self.error:
                    raise self.error
                return self.result

        plain_only = PdfParser(
            plain=StubPdf("plain", "plain result"),  # type: ignore[arg-type]
            ocr=StubPdf("ocr", error=AssertionError("OCR must not run")),  # type: ignore[arg-type]
        )
        self.assertEqual(
            "plain result",
            await plain_only.parse("sample.pdf", enable_ocr=False),
        )
        self.assertEqual(["plain"], calls)

        calls.clear()
        parser = PdfParser(
            plain=StubPdf("plain", {"engine": "pdfreader"}),  # type: ignore[arg-type]
            ocr=StubPdf("ocr", error=PdfParseError("ocr unavailable")),  # type: ignore[arg-type]
        )
        result = await parser.parse("sample.pdf", enable_ocr=True, as_json=True)

        self.assertEqual(["ocr", "plain"], calls)
        self.assertEqual("pdfreader", result["engine"])
        self.assertTrue(result["fallback_used"])

    async def test_plain_pdfreader_handles_small_pdf(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "sample.pdf"
            path.write_bytes(_small_text_pdf())
            parser = PlainPdfParser()
            result = await parser.parse(path)
            outlines = await parser.extract_outlines(path)
        self.assertIn("Hello PDF", result)
        self.assertIsNone(outlines)

    async def test_image_and_doc_external_services_are_stubbed(self) -> None:
        class StubExtractor:
            def __init__(self, result: ExtractorResult) -> None:
                self.result = result

            async def extractor_file(self, *_: object, **__: object) -> ExtractorResult:
                return self.result

        class StubFactory:
            extractors = (
                StubExtractor(
                    ExtractorResult(
                        extractor="first",
                        file_path="image.png",
                        success=False,
                        error="failed",
                    )
                ),
                StubExtractor(
                    ExtractorResult(
                        extractor="second",
                        file_path="image.png",
                        content="识别文本",
                    )
                ),
            )

        image = await ImageParser(factory=StubFactory()).parse(  # type: ignore[arg-type]
            "image.png",
            as_json=True,
        )
        self.assertEqual("second", image["engine"])

        with tempfile.TemporaryDirectory() as directory:
            doc_path = Path(directory) / "sample.doc"
            doc_path.write_bytes(b"fake doc")
            with mock.patch(
                "src.knowledge.parser.doc_parser._request_tika",
                new=mock.AsyncMock(
                    return_value=(
                        "<html><body><p>正文</p></body></html>",
                        "text/html",
                    )
                ),
            ):
                result = await DocParser().parse(doc_path, as_json=True)
        self.assertIn("正文", result["markdown"])


if __name__ == "__main__":
    unittest.main()
